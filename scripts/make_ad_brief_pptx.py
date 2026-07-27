from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

OCEAN = RGBColor(0x00, 0x69, 0x94)
INK = RGBColor(0x0B, 0x2A, 0x38)
MUTED = RGBColor(0x5A, 0x7A, 0x8A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT = RGBColor(0xF0, 0xF7, 0xFA)
LINE = RGBColor(0xD4, 0xE5, 0xEF)


def set_run(run, size=18, bold=False, color=INK):
    run.font.name = "Calibri"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_bg(slide, color=SOFT):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    spTree = slide.shapes._spTree
    sp = shape._element
    spTree.remove(sp)
    spTree.insert(2, sp)


def add_bar(slide):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.18), prs.slide_height
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = OCEAN
    bar.line.fill.background()


def add_text(
    slide,
    left,
    top,
    width,
    height,
    text,
    size=18,
    bold=False,
    color=INK,
    align=PP_ALIGN.LEFT,
):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color)
    return box


def add_bullets(slide, left, top, width, height, items, size=18, color=INK):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.level = 0
        p.space_after = Pt(10)
        run = p.add_run()
        run.text = "•  " + item
        set_run(run, size=size, color=color)
    return box


def footer(slide, page, total=7):
    add_text(
        slide, 0.7, 7.05, 8, 0.3, "Rajeeta  ·  Confidential", size=11, color=MUTED
    )
    add_text(
        slide,
        11.8,
        7.05,
        1.2,
        0.3,
        f"{page} / {total}",
        size=11,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
    )


# Slide 1
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
panel = s.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, 0, Inches(4.6), prs.slide_height
)
panel.fill.solid()
panel.fill.fore_color.rgb = OCEAN
panel.line.fill.background()
add_text(s, 0.55, 2.3, 3.7, 1.2, "RAJEETA", size=40, bold=True, color=WHITE)
add_text(
    s,
    0.55,
    3.3,
    3.7,
    1.2,
    "A short brief for\nour advertising partner",
    size=18,
    color=WHITE,
)
add_text(
    s,
    5.2,
    2.5,
    7.2,
    1.2,
    "What the platform is,\nand what we need from you",
    size=32,
    bold=True,
    color=INK,
)
add_text(
    s,
    5.2,
    4.0,
    7.2,
    0.8,
    "Prepared for campaign planning  ·  Iraq market",
    size=16,
    color=MUTED,
)

# Slide 2
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_bar(s)
add_text(s, 0.7, 0.4, 11, 0.5, "The idea in plain words", size=28, bold=True, color=INK)
add_text(
    s,
    0.7,
    1.2,
    11.5,
    1.2,
    "Rajeeta is a medical platform built for Iraq. Patients book doctors from their phone. Doctors run their clinic from a simple web portal. We sit in the middle and keep appointments, prescriptions, and clinic tools in one place.",
    size=18,
    color=INK,
)
add_bullets(
    s,
    0.7,
    2.9,
    11.5,
    3.5,
    [
        "Patient app (Android) — find a doctor, book a visit, get prescriptions",
        "Doctor website — manage schedule, patients, wallet, and clinic profile",
        "Admin panel — accounts, payments overview, discount codes",
        "Optional AI helper in the app — patient describes symptoms, we suggest a specialty and matching doctors",
    ],
)
footer(s, 2)

# Slide 3
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_bar(s)
add_text(s, 0.7, 0.4, 11, 0.5, "Why this exists", size=28, bold=True, color=INK)
add_text(
    s,
    0.7,
    1.15,
    11.5,
    0.8,
    "Most people still call around or walk in. Clinics lose track of bookings. We want one clean path: open the app, pick a doctor, book, show up.",
    size=18,
    color=INK,
)

cards = [
    (
        "Patients",
        "Less waiting on phone calls.\nClear doctor list by specialty.\nBooking in a few taps.",
    ),
    (
        "Doctors",
        "One place for appointments\nand patient files.\nSimple monthly subscription\nafter a free trial.",
    ),
    (
        "Clinics",
        "Fewer no-shows when\nbookings are confirmed.\nDigital prescriptions\nlinked to the visit.",
    ),
]
for i, (title, body) in enumerate(cards):
    left = 0.7 + i * 4.05
    card = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(2.4),
        Inches(3.8),
        Inches(3.6),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = LINE
    card.adjustments[0] = 0.08
    add_text(s, left + 0.25, 2.65, 3.3, 0.5, title, size=20, bold=True, color=OCEAN)
    add_text(s, left + 0.25, 3.35, 3.3, 2.4, body, size=16, color=INK)
footer(s, 3)

# Slide 4
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_bar(s)
add_text(
    s,
    0.7,
    0.4,
    11,
    0.5,
    "Who the campaign should talk to",
    size=28,
    bold=True,
    color=INK,
)

box = s.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.3), Inches(5.7), Inches(5.0)
)
box.fill.solid()
box.fill.fore_color.rgb = WHITE
box.line.color.rgb = LINE
box.adjustments[0] = 0.06
add_text(s, 1.0, 1.55, 5.1, 0.4, "Primary — Patients", size=20, bold=True, color=OCEAN)
add_bullets(
    s,
    1.0,
    2.2,
    5.1,
    3.8,
    [
        "Ages ~18–45, smartphone users",
        "Baghdad first, then other cities",
        "People who already search for doctors on social / WhatsApp",
        "Parents booking for kids, adults booking themselves",
    ],
    size=16,
)

box2 = s.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.3), Inches(5.7), Inches(5.0)
)
box2.fill.solid()
box2.fill.fore_color.rgb = WHITE
box2.line.color.rgb = LINE
box2.adjustments[0] = 0.06
add_text(s, 7.1, 1.55, 5.1, 0.4, "Secondary — Doctors", size=20, bold=True, color=OCEAN)
add_bullets(
    s,
    7.1,
    2.2,
    5.1,
    3.8,
    [
        "Clinic owners and specialists",
        "Want more booked visits, less chaos",
        "Comfortable with a simple web login",
        "We already have a free trial period to pitch",
    ],
    size=16,
)
footer(s, 4)

# Slide 5
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_bar(s)
add_text(
    s, 0.7, 0.4, 12, 0.5, "What we need from your agency", size=28, bold=True, color=INK
)
add_text(
    s,
    0.7,
    1.05,
    12,
    0.5,
    "Not a full brand rebuild — a launch campaign that gets downloads and doctor sign-ups.",
    size=16,
    color=MUTED,
)
add_bullets(
    s,
    0.7,
    1.8,
    12,
    5,
    [
        "Campaign concept + clear Arabic messaging (Iraqi dialect where it fits; formal Arabic for doctors)",
        "Creative for Meta (Facebook / Instagram) and TikTok — static posts, short videos, stories",
        "Google / YouTube options if budget allows",
        "Landing / store creatives: app icon usage, screenshots layout, short promo video (30–45s)",
        "Media plan: targeting, daily budget split, A/B test ideas for first 30 days",
        "Tracking setup guidance (UTMs, pixel, install events) so we can see what actually works",
        "Weekly report template — installs, cost per install, doctor registrations, top creatives",
    ],
)
footer(s, 5)

# Slide 6
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_bar(s)
add_text(
    s,
    0.7,
    0.4,
    12,
    0.5,
    "Message directions (starting points)",
    size=28,
    bold=True,
    color=INK,
)
add_text(
    s,
    0.7,
    1.1,
    12,
    0.5,
    "You can reshape these — we just want the tone honest and local, not corporate fluff.",
    size=16,
    color=MUTED,
)

rows = [
    ("Patients", "“Book a doctor without the phone chase.”"),
    ("Parents", "“Find the right specialty for your child in minutes.”"),
    ("Doctors", "“Your appointments, in one place — start free.”"),
    ("Trust", "Iraqi platform. Real clinics. Simple tools."),
]
for i, (label, line) in enumerate(rows):
    y = 1.9 + i * 1.1
    tag = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(y), Inches(2.2), Inches(0.75)
    )
    tag.fill.solid()
    tag.fill.fore_color.rgb = OCEAN
    tag.line.fill.background()
    tag.adjustments[0] = 0.2
    add_text(
        s,
        0.85,
        y + 0.18,
        1.9,
        0.45,
        label,
        size=14,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    add_text(s, 3.2, y + 0.15, 9.2, 0.55, line, size=18, color=INK)
footer(s, 6)

# Slide 7
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
add_bar(s)
add_text(s, 0.7, 0.5, 12, 0.5, "Next step", size=28, bold=True, color=INK)
add_text(
    s,
    0.7,
    1.3,
    11.5,
    1.0,
    "We need a proposal from you: concept, sample creatives, estimated media budget for 30–60 days, and what you need from us (logos, screenshots, store links, access).",
    size=18,
    color=INK,
)
add_bullets(
    s,
    0.7,
    2.6,
    11.5,
    3.2,
    [
        "We’ll share: logo pack, app screenshots, doctor portal screenshots, brand colors",
        "We’ll confirm: soft-launch city, install targets, and monthly ad budget",
        "Kickoff call once we agree on scope",
    ],
)
add_text(
    s,
    0.7,
    5.8,
    11.5,
    0.6,
    "Thanks — looking forward to building this with a partner who knows the Iraqi market.",
    size=16,
    color=MUTED,
)
footer(s, 7)

out = r"C:\Users\Mustafa_New\Projects\rajeeta\Rajeeta_Ad_Agency_Brief.pptx"
prs.save(out)
print(out)
